#!/usr/bin/env python3
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def extract_slide_title(slide):
    # Try to find a title placeholder first
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE
            ):
                if hasattr(shape, "text") and shape.text:
                    return shape.text.strip()
        except Exception:
            continue
    # Fallback: first text-bearing shape
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text.strip()
            if txt:
                return txt.splitlines()[0]
    return None

def extract_shape_texts(slide):
    lines = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            tf = shape.text_frame
            for p in tf.paragraphs:
                line = "".join(run.text for run in p.runs).strip()
                if line:
                    lines.append(line)
    return lines

def extract_notes(slide):
    try:
        ns = slide.notes_slide
        if ns and ns.notes_text_frame:
            txt = ns.notes_text_frame.text or ""
            out = [ln.strip() for ln in txt.splitlines() if ln.strip()]
            return out or None
    except Exception:
        pass
    return None

def extract_pptx(path: Path):
    pres = Presentation(str(path))
    slides_out = []
    for i, slide in enumerate(pres.slides, start=1):
        slide_obj = {
            "index": i,
            "title": extract_slide_title(slide),
            "text": extract_shape_texts(slide)
        }
        notes = extract_notes(slide)
        if notes:
            slide_obj["notes"] = notes
        slides_out.append(slide_obj)

    return {
        "file_name": path.name,
        "slide_count": len(slides_out),
        "slides": slides_out
    }

def main():
    ap = argparse.ArgumentParser(description="Extract PPTX text to JSON")
    ap.add_argument("pptx_path", help="Path to .pptx file")
    args = ap.parse_args()

    p = Path(args.pptx_path)
    if not p.exists():
        print(json.dumps({"error": f"File not found: {p}"}), file=sys.stderr)
        sys.exit(2)

    try:
        out = extract_pptx(p)
        # Print JSON to stdout so Node can capture it
        print(json.dumps(out, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
